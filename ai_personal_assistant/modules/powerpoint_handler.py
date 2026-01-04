"""
PowerPoint Automation Module
Handles Microsoft PowerPoint operations
"""

import logging
import os
from typing import Optional, List
from pptx import Presentation
from pptx.util import Inches, Pt
from datetime import datetime

logger = logging.getLogger(__name__)


class PowerPointHandler:
    """Handler for PowerPoint presentation operations"""
    
    def __init__(self, default_template: Optional[str] = None):
        """
        Initialize PowerPoint handler
        
        Args:
            default_template: Path to default template file
        """
        self.default_template = default_template
        self.current_presentation = None
        self.current_file_path = None
        
    def create_presentation(self, title: str = "New Presentation") -> bool:
        """
        Create a new PowerPoint presentation
        
        Args:
            title: Title for the presentation
            
        Returns:
            True if successful, False otherwise
        """
        try:
            if self.default_template and os.path.exists(self.default_template):
                self.current_presentation = Presentation(self.default_template)
            else:
                self.current_presentation = Presentation()
            
            # Add title slide
            title_slide_layout = self.current_presentation.slide_layouts[0]
            slide = self.current_presentation.slides.add_slide(title_slide_layout)
            
            title_shape = slide.shapes.title
            title_shape.text = title
            
            # Add subtitle with date
            if len(slide.placeholders) > 1:
                subtitle = slide.placeholders[1]
                subtitle.text = f"Created on {datetime.now().strftime('%B %d, %Y')}"
            
            logger.info(f"Created new presentation: {title}")
            return True
            
        except Exception as e:
            logger.error(f"Error creating presentation: {e}")
            return False
    
    def add_slide(self, title: str, content: List[str], layout_index: int = 1) -> bool:
        """
        Add a slide to the current presentation
        
        Args:
            title: Slide title
            content: List of content items (bullet points)
            layout_index: Slide layout to use (default: 1 for title and content)
            
        Returns:
            True if successful, False otherwise
        """
        if not self.current_presentation:
            logger.error("No active presentation")
            return False
            
        try:
            slide_layout = self.current_presentation.slide_layouts[layout_index]
            slide = self.current_presentation.slides.add_slide(slide_layout)
            
            # Set title
            title_shape = slide.shapes.title
            title_shape.text = title
            
            # Add content
            if len(slide.placeholders) > 1:
                content_shape = slide.placeholders[1]
                text_frame = content_shape.text_frame
                text_frame.clear()
                
                for item in content:
                    p = text_frame.add_paragraph()
                    p.text = item
                    p.level = 0
            
            logger.info(f"Added slide: {title}")
            return True
            
        except Exception as e:
            logger.error(f"Error adding slide: {e}")
            return False
    
    def save_presentation(self, file_path: Optional[str] = None) -> bool:
        """
        Save the current presentation
        
        Args:
            file_path: Path to save the file (optional)
            
        Returns:
            True if successful, False otherwise
        """
        if not self.current_presentation:
            logger.error("No active presentation to save")
            return False
            
        try:
            if file_path:
                self.current_file_path = file_path
            elif not self.current_file_path:
                # Generate default filename
                timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
                self.current_file_path = f"presentation_{timestamp}.pptx"
            
            self.current_presentation.save(self.current_file_path)
            logger.info(f"Presentation saved: {self.current_file_path}")
            return True
            
        except Exception as e:
            logger.error(f"Error saving presentation: {e}")
            return False
    
    def open_presentation(self, file_path: str) -> bool:
        """
        Open an existing presentation
        
        Args:
            file_path: Path to the presentation file
            
        Returns:
            True if successful, False otherwise
        """
        try:
            if not os.path.exists(file_path):
                logger.error(f"File not found: {file_path}")
                return False
                
            self.current_presentation = Presentation(file_path)
            self.current_file_path = file_path
            logger.info(f"Opened presentation: {file_path}")
            return True
            
        except Exception as e:
            logger.error(f"Error opening presentation: {e}")
            return False
    
    def get_slide_count(self) -> int:
        """
        Get number of slides in current presentation
        
        Returns:
            Number of slides
        """
        if not self.current_presentation:
            return 0
        return len(self.current_presentation.slides)
    
    def close_presentation(self):
        """Close the current presentation"""
        self.current_presentation = None
        self.current_file_path = None
        logger.info("Closed current presentation")
